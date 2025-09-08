import os
import json
from flask import Flask, request, jsonify, render_template, send_from_directory, send_file
from werkzeug.utils import secure_filename
from pptx import Presentation
import openpyxl
import requests
import PyPDF2
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from dotenv import load_dotenv
import dashscope
from http import HTTPStatus

load_dotenv()
# API密钥现在通过手动配置管理，不再使用硬编码密钥

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['ALLOWED_EXTENSIONS'] = {'pptx', 'pdf', 'xlsx', 'mp3', 'wav'}

if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

def allowed_file(filename, allowed_extensions=None):
    if allowed_extensions is None:
        allowed_extensions = app.config['ALLOWED_EXTENSIONS']
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_extensions

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/static/<path:path>')
def serve_static(path):
    return send_from_directory('static', path)

# 上传文件接口
@app.route('/api/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        return jsonify({'filename': filename, 'file_path': file_path}), 200
    return jsonify({'error': 'File type not allowed'}), 400

# 解析文档内容（支持PPT和PDF）
@app.route('/api/parse_document', methods=['POST'])
def parse_document():
    data = request.json
    doc_path = data.get('doc_path')
    if not doc_path or not os.path.exists(doc_path):
        return jsonify({'error': '文件未找到'}), 400
    
    try:
        if doc_path.lower().endswith('.pptx'):
            # 解析PPT
            prs = Presentation(doc_path)
            doc_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        doc_text.append(shape.text)
            return jsonify({'doc_text': '\n'.join(doc_text)}), 200
            
        elif doc_path.lower().endswith('.pdf'):
            # 解析PDF
            pdf_text = []
            with open(doc_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page_num in range(len(reader.pages)):
                    page = reader.pages[page_num]
                    pdf_text.append(page.extract_text())
            return jsonify({'doc_text': '\n'.join(pdf_text)}), 200
            
        else:
            return jsonify({'error': '不支持的文件格式'}), 400
    except Exception as e:
        return jsonify({'error': f'解析文件错误: {str(e)}'}), 500

# 保留旧接口以兼容
@app.route('/api/parse_ppt', methods=['POST'])
def parse_ppt():
    data = request.json
    ppt_path = data.get('ppt_path')
    if not ppt_path or not os.path.exists(ppt_path):
        return jsonify({'error': 'PPT file not found'}), 400
    
    # 将请求转发给新接口
    result = parse_document()
    
    # 转换响应格式以保持兼容性
    if result[1] == 200:
        response_data = result[0].get_json()
        return jsonify({'ppt_text': response_data.get('doc_text', '')}), 200
    return result

# 解析评分表内容
@app.route('/api/parse_score', methods=['POST'])
def parse_score():
    data = request.json
    score_path = data.get('score_path')
    if not score_path or not os.path.exists(score_path):
        return jsonify({'error': 'Score file not found'}), 400
    try:
        wb = openpyxl.load_workbook(score_path)
        ws = wb.active
        score_items = []
        for row in ws.iter_rows(min_row=2, values_only=True):
            if row and any(row):  # 确保行不为空
                score_items.append(' '.join(str(cell) for cell in row if cell))
        return jsonify({'score_items': score_items}), 200
    except Exception as e:
        return jsonify({'error': f'Error parsing score file: {str(e)}'}), 500

# 文档与评分表内容一致性校对
@app.route('/api/verify', methods=['POST'])
def verify():
    data = request.json
    doc_text = data.get('doc_text', '') or data.get('ppt_text', '')  # 兼容旧版本
    score_items = data.get('score_items', [])
    
    # 调用大模型进行智能分析，而不是简单的文本匹配
    analysis_result = call_dashscope_llm_analysis(doc_text, score_items)
    
    return jsonify({
        'missing_items': analysis_result.get('missing_items', []), 
        'suggestions': analysis_result.get('suggestions', ''),
        'success': True
    }), 200

# 调用阿里云DashScope大模型进行智能分析
# API密钥现在通过手动配置管理，不再使用硬编码密钥

def call_dashscope_llm_analysis(doc_text, score_items):
    """
    使用专业的述职报告校对助手进行智能分析
    """
    prompt = f"""
## 角色
你是一个专业的述职报告校对助手。你的任务是严格、细致地对比"评价表"中的要求和"述职PPT"中的实际内容，找出所有不一致、缺失或表述不清的地方，并提供具体、可行的修改参考。

## 任务目标
根据我提供的[评价表内容]和[PPT内容]，完成以下任务：

识别差异：找出PPT中未能满足评价表要求的所有项。
描述问题：清晰地描述每一处不一致或缺失的具体情况。
提供建议：针对每一个问题，给出建设性的、可操作的修改参考。

## 输入
[评价表内容]: 
{chr(10).join([f"- {item}" for item in score_items if item.strip()])}

[PPT内容]: 
{doc_text}

## Workflow
请严格按照以下步骤执行：
1. 理解要求：首先，仔细阅读并完全理解[评价表内容]中的每一条标准和要求。将其视为一个必须逐项检查的清单（Checklist）。
2. 内容匹配：然后，通读[PPT内容]，将PPT的每一部分与评价表中的要求进行逐一匹配和核对。
3. 发现问题：识别出以下几类问题：
-完全缺失：评价表要求的内容，在PPT中完全没有提及。
-部分缺失/不充分：PPT中虽然提及了相关内容，但不够深入、具体，未能完全满足评价表的要求（例如，要求数据对比，但只给了结论）。
-表述不符：PPT的表述与评价表的要求有出入。
4. 生成结果：针对每一个发现的问题，生成一个包含"问题描述"和"修改参考"的组合。

## 输出要求
请严格按照以下JSON格式返回，不要添加任何其他文字：

```json
{{
    "missing_items": ["问题1的简要描述", "问题2的简要描述"],
    "suggestions": "HTML格式的详细修改参考"
}}
```

suggestions字段必须按以下HTML格式输出，直接列出修改参考：
```html
<div class="suggestions-content">
<ul>
<li>具体的修改参考1</li>
<li>具体的修改参考2</li>
<li>具体的修改参考3</li>
</ul>
</div>
```

注意：直接列出修改参考即可，每条用一句话简洁描述，不需要问题分类。

## Rules
1. 必须严格返回JSON格式，开头是{{，结尾是}}
2. 不要添加```json```代码块标记
3. 不要添加任何解释性文字
4. missing_items如果没有问题就返回空数组[]
5. 保持输出简洁，每个问题和修改参考都用一句话描述，避免冗长内容

示例输出：
{{"missing_items": ["工作成果量化不足"], "suggestions": "<div class=\"suggestions-content\"><h5>系统已完成校对，共检测到 1处 需要关注的内容</h5><h5>一、工作与能力展示维度</h5><ul><li><strong>问题：工作成果量化不足</strong><br/>修改参考：补充具体数据和案例</li></ul></div>"}}
    """
    
    url = 'https://dashscope.aliyuncs.com/api/v1/services/aigc/text-generation/generation'
    headers = {
        'Authorization': f'Bearer {DASHSCOPE_API_KEY}',
        'Content-Type': 'application/json'
    }
    payload = {
        'model': 'qwen-plus',
        'input': {'prompt': prompt},
        'parameters': {'result_format': 'text', 'max_tokens': 2000}
    }
    
    try:
        resp = requests.post(url, headers=headers, json=payload, timeout=60)
        if resp.status_code == 200:
            result = resp.json()
            response_text = result.get('output', {}).get('text', '')
            
            # 尝试解析JSON响应
            try:
                # 清理可能的代码块标记
                cleaned_text = response_text.strip()
                if cleaned_text.startswith('```json'):
                    cleaned_text = cleaned_text[7:]
                if cleaned_text.endswith('```'):
                    cleaned_text = cleaned_text[:-3]
                cleaned_text = cleaned_text.strip()
                
                analysis_result = json.loads(cleaned_text)
                return analysis_result
            except json.JSONDecodeError as e:
                # 如果JSON解析失败，返回基础分析结果，包含调试信息
                return {
                    'missing_items': ['大模型返回格式解析失败'],
                    'suggestions': f'''
                    <div class="suggestions-content">
                    <div class="alert alert-warning">
                        <h6><i class="fas fa-exclamation-triangle me-2"></i>JSON解析错误</h6>
                        <p><strong>错误位置：</strong> 第{e.lineno}行，第{e.colno}列</p>
                        <p><strong>错误信息：</strong> {e.msg}</p>
                        <details>
                            <summary>查看原始响应</summary>
                            <pre style="max-height: 200px; overflow-y: auto; font-size: 0.8em;">{response_text[:1000]}{'...' if len(response_text) > 1000 else ''}</pre>
                        </details>
                    </div>
                    <h5>修改参考</h5>
                    <ul>
                        <li>大模型返回格式错误，建议稍后重试或检查网络连接</li>
                    </ul>
                    </div>
                    '''
                }
        else:
            return get_fallback_analysis(doc_text, score_items)
    except Exception as e:
        return get_fallback_analysis(doc_text, score_items, str(e))

def get_fallback_analysis(doc_text, score_items, error_msg=None):
    """备用分析逻辑，遵循专业校对格式"""
    missing_items = []
    
    # 简单关键词匹配分析
    keywords_to_check = {
        '工作成果展示': ['成果', '业绩', '结果', '产出', '完成'],
        '团队协作能力': ['协作', '团队', '合作', '配合', '沟通'],
        '能力发展体现': ['学习', '成长', '提升', '发展', '进步'],
        '业务影响说明': ['业务', '价值', '效益', '收益', '贡献'],
        '创新思维展现': ['创新', '改进', '优化', '新方法', '突破']
    }
    
    detected_issues = []
    for category, keywords in keywords_to_check.items():
        if not any(keyword in doc_text for keyword in keywords):
            missing_items.append(f"{category}内容不充分")
            detected_issues.append(category)
    
    # 生成专业格式的建议
    issue_count = len(detected_issues)
    mode_tip = "⚠️ 备用分析模式" if error_msg else "⚠️ 基础分析模式"
    suggestions_content = f"""
    <div class="suggestions-content">
    <div class="alert alert-warning">
        <i class="fas fa-exclamation-triangle me-2"></i>
        <strong>{mode_tip}</strong> - 智能分析服务不可用，当前使用基础关键词分析，建议稍后重试获得更专业的校对结果。
    </div>
    <h5>系统已完成校对，共检测到 {issue_count}处 需要关注的内容</h5>
    
    <h5>一、内容完整性维度</h5>
    <ul>
    """
    
    if '工作成果展示' in detected_issues:
        suggestions_content += """
        <li><strong>问题：工作成果数据不够具体</strong><br/>
        修改参考：补充项目数量、质量指标等具体数据和案例</li>
        """
    
    if '团队协作能力' in detected_issues:
        suggestions_content += """
        <li><strong>问题：团队协作案例缺失</strong><br/>
        修改参考：添加跨部门协作具体案例和个人贡献描述</li>
        """
    
    if '业务影响说明' in detected_issues:
        suggestions_content += """
        <li><strong>问题：业务价值体现不足</strong><br/>
        修改参考：量化工作对业务的具体贡献和价值影响</li>
        """
    
    suggestions_content += """
    </ul>
    
    <h5>二、表达与呈现维度</h5>
    <ul>
        <li><strong>问题：分析基于基础模式</strong><br/>
    """
    
    if error_msg:
        suggestions_content += f"修改参考：建议检查网络连接后重新校对获得更精准分析</li>"
    else:
        suggestions_content += "修改参考：建议稍后重试获得更专业的校对结果</li>"
    
    suggestions_content += """
    </ul>
    </div>
    """
    
    return {
        'missing_items': missing_items,
        'suggestions': suggestions_content
    }

# 保留原有函数以兼容其他功能
def call_dashscope_llm(doc_text, score_items, missing_items):
    prompt = f"""
    任务：分析述职文档与评分表的一致性，并提供针对性建议
    
    文档内容概述：
    {doc_text[:500]}...（内容省略）
    
    评分表考核项：
    {', '.join(score_items[:10])}...（可能有更多项）
    
    缺失项：
    {', '.join(missing_items) if missing_items else '无缺失项'}
    
    请根据上述信息，提供详细的修改参考，严格按照以下HTML格式输出：
    
    <div class="suggestions-content">
    <h5>一、如何补充缺失的评分项内容</h5>
    <ul>
    <li>具体建议1</li>
    <li>具体建议2</li>
    </ul>
    
    <h5>二、如何更好地将文档内容与评分表关键指标对齐</h5>
    <ul>
    <li>具体建议1</li>
    <li>具体建议2</li>
    </ul>
    
    <h5>三、提升述职效果的具体建议</h5>
    <ul>
    <li>具体建议1</li>
    <li>具体建议2</li>
    </ul>
    </div>
    
    请严格按照上述HTML结构输出，不要添加其他文字说明。
    """
    
    url = 'https://dashscope.aliyuncs.com/api/v1/services/aigc/text-generation/generation'
    headers = {
        'Authorization': f'Bearer {DASHSCOPE_API_KEY}',
        'Content-Type': 'application/json'
    }
    payload = {
        'model': 'qwen-plus',
        'input': {'prompt': prompt},
        'parameters': {'result_format': 'text'}
    }
    
    try:
        resp = requests.post(url, headers=headers, json=payload, timeout=30)
        if resp.status_code == 200:
            result = resp.json()
            return result.get('output', {}).get('text', '')
        else:
            # 模拟返回结果，API调用失败时使用
            if missing_items:
                return """
                <div class="suggestions-content">
                <h5>一、如何补充缺失的评分项内容</h5>
                <ul>
                    <li>针对每个缺失项，在文档中添加专门的章节或段落来详细阐述</li>
                    <li>提供具体的量化数据和实际案例来支撑相关内容</li>
                    <li>确保覆盖评分表中的所有关键考核维度</li>
                </ul>
                
                <h5>二、如何更好地将文档内容与评分表关键指标对齐</h5>
                <ul>
                    <li>在文档开头添加目录，与评分表各维度一一对应</li>
                    <li>为每个评分项提供清晰的标题和结构化内容</li>
                    <li>使用关键词和标签帮助评委快速定位相关内容</li>
                </ul>
                
                <h5>三、提升述职效果的具体建议</h5>
                <ul>
                    <li>突出个人贡献，明确区分团队成就与个人作用</li>
                    <li>添加项目前后对比数据，展示实际效果和价值</li>
                    <li>使用图表和可视化元素提升表达效果</li>
                    <li>确保内容简洁明了，重点突出</li>
                </ul>
                </div>
                """
            else:
                return """
                <div class="suggestions-content">
                <h5>文档内容完整性分析</h5>
                <ul>
                    <li>您的文档内容已经全面覆盖了评分表的所有要求，保持这种全面性</li>
                    <li>继续为每个评分项提供充分的量化数据和具体案例支持</li>
                </ul>
                
                <h5>进一步优化建议</h5>
                <ul>
                    <li>使用更简洁明了的语言表达，提高可读性</li>
                    <li>考虑添加图表可视化展示关键成就</li>
                    <li>确保重点内容突出，层次分明</li>
                </ul>
                </div>
                """
    except Exception as e:
        # 模拟返回结果，发生异常时使用
        return f"""
        <div class="suggestions-content">
        <h5>系统提示</h5>
        <ul>
            <li>大模型调用暂时不可用 ({str(e)})</li>
            <li>已启用本地分析模式，提供基础建议</li>
        </ul>
        
        <h5>通用优化建议</h5>
        <ul>
            <li>确保文档内容与评分表各项条款一一对应</li>
            <li>使用具体数据和案例支持您的陈述</li>
            <li>突出个人贡献和核心价值</li>
            <li>保持内容结构清晰，重点明确</li>
        </ul>
        </div>
        """

# 员工诊断报告API
def extract_employee_info_with_llm(text, filename):
    """使用大模型分析文档内容并提取员工信息"""
    prompt = f"""
## 任务
请分析以下文档内容，提取员工的基本信息。

## 文档文件名
{filename}

## 文档内容（前2000字符）
{text[:2000]}

## 输出要求
请严格按照以下格式输出员工信息，如果某项信息无法确定，请填写"未知"：

员工姓名：[姓名]
职位信息：[职位/岗位/职级]
评估周期：[年份和季度]

## 注意事项
1. 姓名通常出现在：述职人、姓名、汇报人、花名、申请人等字段后
2. 职位信息可能包括：职位、岗位、职级、申请岗位等
3. 评估周期可能出现在文件名或文档内容中，格式如：2025年第三季度、2025Q3等
4. 严格按照指定格式输出，不要添加其他说明文字
"""
    
    try:
        print(f"[DEBUG] 开始使用大模型提取员工信息，文件名: {filename}")
        print(f"[DEBUG] 文档内容前500字符: {text[:500]}")
        
        response = dashscope.Generation.call(
            model='qwen-plus',
            prompt=prompt,
            result_format='text',
            api_key=DASHSCOPE_API_KEY
        )
        
        if response.status_code == HTTPStatus.OK:
            response_text = response.output.text.strip()
            print(f"[DEBUG] 大模型返回结果: {response_text}")
            
            # 使用正则表达式从大模型输出中提取信息
            import re
            
            employee_name = '未知'
            position = '未知'
            quarter = '未知'
            
            # 提取员工姓名
            name_match = re.search(r'员工姓名[：:]\s*([^\n\r]+)', response_text)
            if name_match:
                employee_name = name_match.group(1).strip()
                print(f"[DEBUG] 提取到员工姓名: {employee_name}")
            
            # 提取职位信息
            position_match = re.search(r'职位信息[：:]\s*([^\n\r]+)', response_text)
            if position_match:
                position = position_match.group(1).strip()
                print(f"[DEBUG] 提取到职位信息: {position}")
            
            # 提取评估周期
            quarter_match = re.search(r'评估周期[：:]\s*([^\n\r]+)', response_text)
            if quarter_match:
                quarter = quarter_match.group(1).strip()
                print(f"[DEBUG] 提取到评估周期: {quarter}")
            
            print(f"[DEBUG] 最终提取结果 - 姓名: {employee_name}, 职位: {position}, 周期: {quarter}")
            return employee_name, position, quarter
        else:
            print(f"[DEBUG] 大模型调用失败: {response.status_code}, {response.message}")
            
    except Exception as e:
        print(f"大模型提取员工信息失败: {str(e)}")
    
    # 如果大模型提取失败，回退到原有的正则表达式方法
    print(f"[DEBUG] 回退到正则提取方法")
    return extract_employee_info_from_content_fallback(text, filename)

def extract_employee_info_from_content_fallback(text, filename):
    """备用的正则表达式提取方法"""
    import re
    
    employee_name = '未知'
    position = '未知'
    quarter = '未知'
    
    # 从文件名提取信息
    filename_lower = filename.lower()
    
    # 提取季度信息
    quarter_patterns = [
        r'(\d{4})年?第?([一二三四1234])季度?',
        r'(\d{4})年?([一二三四1234])季度?',
        r'(\d{4})q([1234])',
        r'q([1234])(\d{4})',
        r'(\d{4})[年-](\d{1,2})月?'
    ]
    
    for pattern in quarter_patterns:
        match = re.search(pattern, filename_lower + text[:500])
        if match:
            if '月' in pattern:
                year, month = match.groups()
                month_num = int(month)
                if 1 <= month_num <= 3:
                    quarter = f'{year}年第一季度'
                elif 4 <= month_num <= 6:
                    quarter = f'{year}年第二季度'
                elif 7 <= month_num <= 9:
                    quarter = f'{year}年第三季度'
                elif 10 <= month_num <= 12:
                    quarter = f'{year}年第四季度'
            else:
                groups = match.groups()
                if len(groups) >= 2:
                    year = groups[0] if groups[0].isdigit() else groups[1]
                    q = groups[1] if groups[0].isdigit() else groups[0]
                    quarter_map = {'1': '一', '2': '二', '3': '三', '4': '四',
                                 '一': '一', '二': '二', '三': '三', '四': '四'}
                    if q in quarter_map:
                        quarter = f'{year}年第{quarter_map[q]}季度'
            break
    
    # 从文本内容提取员工姓名
    name_patterns = [
        r'述职人[：:]\s*([\u4e00-\u9fa5]{2,4})',
        r'姓名[：:]\s*([\u4e00-\u9fa5]{2,4})',
        r'汇报人[：:]\s*([\u4e00-\u9fa5]{2,4})',
        r'花名[：:]\s*([\u4e00-\u9fa5]{2,4})',
        r'我是\s*([\u4e00-\u9fa5]{2,4})',
        r'本人\s*([\u4e00-\u9fa5]{2,4})',
        r'申请人[：:]\s*([\u4e00-\u9fa5]{2,4})'
    ]
    
    for pattern in name_patterns:
        match = re.search(pattern, text[:1500])
        if match:
            employee_name = match.group(1).strip()
            break
    
    # 从文本内容提取职位信息
    position_patterns = [
        r'职位[：:]\s*([\u4e00-\u9fa5]{2,10})',
        r'岗位[：:]\s*([\u4e00-\u9fa5]{2,10})',
        r'职级[：:]\s*(p\d+|P\d+|[\u4e00-\u9fa5]{2,10})',
        r'申请岗位[：:]\s*([\u4e00-\u9fa5]{2,10})',
        r'担任\s*([\u4e00-\u9fa5]{2,10})',
        r'(高级|中级|初级|资深)?(工程师|开发|架构师|经理|主管|总监|专员|策划)',
        r'(P\d+|T\d+|M\d+|p\d+|t\d+|m\d+)'
    ]
    
    for pattern in position_patterns:
        match = re.search(pattern, text[:1500], re.IGNORECASE)
        if match:
            if match.groups():
                position = ''.join(filter(None, match.groups())).strip()
            else:
                position = match.group(0).strip()
            break
    
    return employee_name, position, quarter

def extract_employee_info_from_content(text, filename):
    """从文件内容和文件名中提取员工信息（主入口函数）"""
    return extract_employee_info_with_llm(text, filename)

@app.route('/api/generate_diagnosis', methods=['POST'])
def generate_diagnosis():
    data = request.get_json()
    
    # 获取请求参数
    employee_name = data.get('employee_name', '未知员工')
    ability_model = data.get('ability_model', '通用能力模型')
    quarter = data.get('quarter', '未知季度')
    doc_path = data.get('doc_path')
    # 不再使用音频文件
    # audio_path = data.get('audio_path', '')
    
    if not doc_path:
        return jsonify({'error': 'Missing doc_path'}), 400
    
    try:
        # 读取文档内容
        if doc_path.startswith('uploads/') or doc_path.startswith('uploads\\'):
            full_doc_path = doc_path
        else:
            full_doc_path = os.path.join(app.config['UPLOAD_FOLDER'], doc_path)
            
        report_text = ''
        filename = os.path.basename(doc_path)
        
        if doc_path.endswith('.pdf'):
            with open(full_doc_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    report_text += page.extract_text() + '\n'
        elif doc_path.endswith('.pptx'):
            prs = Presentation(full_doc_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        report_text += shape.text + '\n'
        else:
            return jsonify({'error': 'Unsupported file format'}), 400
            
        if not report_text.strip():
            return jsonify({'error': 'Failed to extract text from document'}), 400
            
        # 从文档内容和文件名中提取员工信息
        extracted_name, extracted_position, extracted_quarter = extract_employee_info_from_content(report_text, filename)
        
        # 如果前端没有提供员工信息，则使用提取的信息
        if employee_name == '未知员工' or not employee_name:
            employee_name = extracted_name
        if quarter == '未知季度' or not quarter:
            quarter = extracted_quarter
            
        # 调用千问模型进行诊断分析
        diagnosis_result = call_qianwen_for_diagnosis(report_text, employee_name, ability_model, quarter)
        
        if diagnosis_result:
            # 更新诊断结果中的员工信息
            diagnosis_result['employee_info']['name'] = employee_name
            diagnosis_result['employee_info']['quarter'] = quarter
            if extracted_position != '未知':
                diagnosis_result['employee_info']['position'] = extracted_position
            
            return jsonify({
                'success': True,
                'diagnosis': diagnosis_result,
                'extracted_info': {
                    'name': extracted_name,
                    'position': extracted_position,
                    'quarter': extracted_quarter
                }
            }), 200
        else:
            # 如果千问模型调用失败，返回备用诊断结果
            fallback_result = get_fallback_diagnosis(employee_name, ability_model, quarter)
            # 更新备用结果中的员工信息
            fallback_result['employee_info']['name'] = employee_name
            fallback_result['employee_info']['quarter'] = quarter
            if extracted_position != '未知':
                fallback_result['employee_info']['position'] = extracted_position
                
            return jsonify({
                'success': True,
                'diagnosis': fallback_result,
                'note': '使用备用分析模式',
                'extracted_info': {
                    'name': extracted_name,
                    'position': extracted_position,
                    'quarter': extracted_quarter
                }
            }), 200
            
    except FileNotFoundError:
        return jsonify({'error': f'Failed to read document: File not found: {doc_path}'}), 400
    except Exception as e:
        return jsonify({'error': f'Failed to generate diagnosis: {str(e)}'}), 500

# 群体分析API
@app.route('/api/generate_cohort_analysis', methods=['POST'])
def generate_cohort_analysis():
    # 这里可以实现真实的群体分析报告生成逻辑
    # 当前返回模拟数据
    return jsonify({
        'success': True,
        'cohort_analysis': {
            'cohort_name': '所有P6级员工',
            'model_name': '技术序列P6能力模型',
            'time_range': '2023年全年',
            'average_abilities': {
                'technical_innovation': 4.2,
                'business_impact': 3.4,
                'teamwork': 4.3, 
                'project_management': 4.0,
                'cost_awareness': 3.1,
                'strategic_thinking': 3.5
            },
            'strengths': [
                {'name': '技术创新', 'score': 4.2, 'description': '普遍表现突出，能够快速学习新技术'},
                {'name': '团队协作', 'score': 4.3, 'description': '协作意识强，能有效促进团队合作'},
                {'name': '执行力', 'score': 4.1, 'description': '任务完成度高，交付质量好'}
            ],
            'weaknesses': [
                {'name': '成本意识', 'score': 3.1, 'description': '普遍缺乏成本效益考量'},
                {'name': '业务影响力', 'score': 3.4, 'description': '技术工作与业务价值关联不够明确'},
                {'name': '战略思维', 'score': 3.5, 'description': '全局观和长远规划能力有待加强'}
            ],
            'best_practices': [
                {
                    'ability': '业务影响力',
                    'employee': '员工A',
                    'score': 4.8,
                    'practice': '在项目启动前主动与业务团队沟通，明确业务目标和预期价值，并在项目过程中持续对齐',
                    'result': '项目最终用户采纳率超过85%，直接贡献营收增长12%'
                },
                {
                    'ability': '业务影响力',
                    'employee': '员工B',
                    'score': 4.6,
                    'practice': '建立了技术工作业务价值评估框架，对所有开发任务进行价值评分和优先级排序',
                    'result': '团队资源分配效率提升30%，高价值任务完成量增加40%'
                }
            ]
        }
    }), 200

# 打分建议API
@app.route('/api/generate_scoring_suggestion', methods=['POST'])
def generate_scoring_suggestion():
    data = request.get_json()
    
    # 支持两种数据格式：新格式（前端发送的）和旧格式（直接传递文本）
    if 'report_text' in data and 'scoring_table_text' in data:
        # 旧格式：直接使用传递的文本
        report_text = data.get('report_text')
        scoring_table_text = data.get('scoring_table_text')
    else:
        # 新格式：从前端参数中构建
        employee_name = data.get('employee_name')
        ability_model = data.get('ability_model')
        quarter = data.get('quarter')
        doc_path = data.get('doc_path')
        
        if not doc_path:
            return jsonify({"error": "Missing doc_path"}), 400
            
        # 从文档路径读取文本内容
        try:
            # 检查 doc_path 是否已经包含 uploads 前缀
            if doc_path.startswith('uploads/') or doc_path.startswith('uploads\\'):
                full_doc_path = doc_path
            else:
                full_doc_path = os.path.join(app.config['UPLOAD_FOLDER'], doc_path)
            if doc_path.endswith('.pdf'):
                with open(full_doc_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    report_text = ''
                    for page in reader.pages:
                        report_text += page.extract_text() + '\n'
            elif doc_path.endswith('.pptx'):
                prs = Presentation(full_doc_path)
                report_text = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            report_text += shape.text + '\n'
            else:
                return jsonify({"error": "Unsupported file format"}), 400
                
            if not report_text.strip():
                return jsonify({"error": "No text content found in document"}), 400
                
        except Exception as e:
            return jsonify({"error": f"Failed to read document: {str(e)}"}), 500
            
        # 构建评分表文本（基于能力模型）
        scoring_table_text = f"员工：{employee_name}\n职位：{ability_model}\n季度：{quarter}\n\n能力评估维度：\n1. 技术掌握与应用\n2. 项目管理能力\n3. 团队协作与沟通\n4. 业务理解与贡献\n5. 学习能力与创新思维"

    if not report_text or not scoring_table_text:
        return jsonify({"error": "Missing report_text or scoring_table_text"}), 400

    analysis_result = call_bailian_for_suggestion(report_text, scoring_table_text)
    
    if not analysis_result:
        analysis_result = get_fallback_suggestion(report_text, scoring_table_text)

    if analysis_result:
        # 为了兼容前端，包装返回结果
        return jsonify({
            "success": True,
            "scoring": {
                "employee_name": data.get('employee_name', '未知员工'),
                "position": data.get('ability_model', '未知职位'),
                "quarter": data.get('quarter', '未知季度'),
                "core_strengths": analysis_result.get('core_strengths', ''),
                "areas_for_development": analysis_result.get('areas_for_development', ''),
                "scoring_suggestions": analysis_result.get('scoring_suggestions', []),
                "abilities": [item.get('ability', '') for item in analysis_result.get('scoring_suggestions', [])]
            }
        })
    else:
        return jsonify({"error": "Failed to get analysis from LLM and fallback"}), 500

def call_bailian_for_suggestion(report_text, scoring_table_text):
    """
    Calls the Alibaba Cloud Bailian LLM to get scoring suggestions.
    """
    system_prompt = """
# Role: 资深人力资源专家和绩效评估顾问
你是一位拥有15年以上经验的资深HR专家，专精于员工述职报告分析和绩效评估。你具备敏锐的洞察力，能够从述职报告中准确识别员工的核心能力、工作亮点和发展潜力，并提供精准、可操作的改进建议。

# Context: 
你正在分析一份员工述职报告，需要根据公司的评分体系进行客观、专业的评估。你的分析将直接影响员工的绩效评定和职业发展规划，因此必须确保评估的准确性和建设性。

# 评分标准（严格遵循）:
- **5分（卓越）**: 表现远超预期，在该能力项上有突出贡献或创新突破，可作为标杆
- **4分（优秀）**: 表现超出预期，展现了较强的该项能力，有明显亮点
- **3分（良好）**: 表现符合预期，基本达到岗位要求，无明显不足
- **2分（待改进）**: 表现略低于预期，该能力项需要重点关注和改进
- **1分（不足）**: 表现明显低于预期，该能力项存在明显缺陷，需要紧急提升

# 分析要求:
1. **证据导向**: 每个评分必须基于报告中的具体事实、数据或案例
2. **客观公正**: 避免主观臆断，基于事实进行评估
3. **建设性**: 提供具体、可操作的改进建议
4. **全面性**: 既要肯定优势，也要指出不足
5. **专业性**: 使用专业的HR术语和评估方法

# 输入内容:
**述职报告内容:**
{report_text}

**评分体系:**
{scoring_table_text}

# 分析示例（参考格式和深度）:
```json
{{
  "core_strengths": "技术创新能力突出，项目交付质量高，团队协作意识强，具备优秀的问题解决能力",
  "areas_for_development": "战略思维和前瞻性规划需要加强，跨部门沟通协调能力有待提升",
  "scoring_suggestions": [
    {{
      "ability": "技术掌握与应用",
      "score": 5,
      "basis": "报告中明确提到'独立设计并实现了核心算法优化方案，系统性能提升40%，获得公司技术创新奖'，展现了卓越的技术能力和创新思维，成果具有重大价值",
      "suggestion": "建议将技术创新经验进行系统化总结，通过内部技术分享会或技术博客的形式，带动团队整体技术水平提升，扩大个人技术影响力"
    }},
    {{
      "ability": "项目管理能力",
      "score": 4,
      "basis": "报告显示'同时负责3个重点项目，均按时保质交付，客户满意度达到95%以上，项目成本控制在预算范围内'，体现了较强的项目统筹和执行能力",
      "suggestion": "建议系统学习PMP或敏捷项目管理方法论，进一步提升复杂项目的风险识别和应对能力，可考虑承担更大规模的跨部门项目"
    }}
  ]
}}
```

# 输出要求:
请严格按照以下JSON格式输出，确保:
1. 格式完全正确，可直接解析
2. 不包含任何markdown标记或额外文本
3. 评分依据必须引用具体的报告内容
4. 建议必须具体可操作

{{
  "core_strengths": "（核心优势总结，35-50字，要具体且有针对性）",
  "areas_for_development": "（待发展领域总结，35-50字，要建设性且可操作）",
  "scoring_suggestions": [
    {{
      "ability": "（能力项名称，与评分体系完全一致）",
      "score": (1-5的整数评分),
      "basis": "（评分依据，必须引用报告具体内容，120-180字）",
      "suggestion": "（具体可操作的提升建议，100-150字）"
    }}
  ]
}}
"""
    messages = [
        {'role': 'system', 'content': 'You are a helpful assistant.'},
        {'role': 'user', 'content': system_prompt.format(report_text=report_text, scoring_table_text=scoring_table_text)}
    ]

    try:
        response = dashscope.Generation.call(
            model='qwen-max',
            messages=messages,
            result_format='message',
        )
        if response.status_code == HTTPStatus.OK:
            content = response.output.choices[0].message.content
            # The response might be wrapped in ```json ... ```, so we need to extract it.
            if content.strip().startswith("```json"):
                content = content.strip()[7:-3]
            return json.loads(content)
        else:
            print(f"Error from Bailian API: {response.code} - {response.message}")
            return None
    except Exception as e:
        print(f"Exception while calling Bailian API: {e}")
        return None

def get_fallback_suggestion(report_text, scoring_table_text):
    """
    如果LLM调用失败，返回一个备用的、写死的分析结果。
    """
    print("LLM调用失败，使用备用分析结果。")
    # 这里返回一个结构与LLM输出一致的写死JSON对象
    return {
        "core_strengths": "这是一个备用的核心优势分析。在报告中，您展现了出色的项目执行能力和团队沟通技巧。",
        "areas_for_development": "这是一个备用的待发展领域分析。建议在战略规划和跨部门协作方面投入更多精力。",
        "scoring_suggestions": [
            {
                "ability": "技术掌握与应用",
                "score": 4,
                "basis": "备用依据：报告中提到了您在XX项目中成功应用了新技术，解决了关键问题。",
                "suggestion": "备用建议：持续学习前沿技术，并尝试在团队内部进行分享，扩大技术影响力。"
            },
            {
                "ability": "团队协作与沟通",
                "score": 5,
                "basis": "备用依据：报告中多次提到您与团队成员紧密合作，共同完成了挑战性任务。",
                "suggestion": "备用建议：您的沟通协作能力已经很强，可以尝试承担更复杂的跨团队沟通角色。"
            },
            {
                "ability": "业务理解与贡献",
                "score": 3,
                "basis": "备用依据：报告中对业务的理解较为到位，但对业务的贡献描述不够具体。",
                "suggestion": "备用建议：在未来的报告中，多使用可量化的数据来展示您对业务的具体贡献。"
            }
        ]
    }

def call_qianwen_for_diagnosis(report_text, employee_name, ability_model, quarter):
    """
    调用千问模型进行员工个人诊断分析
    """
    prompt = f"""
## 角色
你是一位资深的人力资源专家和职业发展顾问，具有丰富的员工能力评估和职业发展指导经验。

## 任务目标
基于员工的述职报告内容，进行全面的个人诊断分析，包括能力评估、优势识别、发展领域分析和具体的成长建议。

## 输入信息
**员工姓名**: {employee_name}
**能力模型**: {ability_model}
**评估周期**: {quarter}
**述职报告内容**:
{report_text}

## 分析维度
请从以下六个核心维度进行分析评估（1-5分制）：
1. **技术创新能力** - 学习新技术、创新解决方案的能力
2. **业务影响力** - 对业务目标达成的贡献和影响
3. **团队协作能力** - 沟通协调、团队合作的表现
4. **项目管理能力** - 项目规划、执行、交付的能力
5. **成本意识** - 资源优化、成本控制的意识和行为
6. **战略思维** - 全局视角、长远规划的思考能力

## 输出要求
请严格按照以下JSON格式输出分析结果：

{{
    "employee_info": {{
        "name": "{employee_name}",
        "position": "{ability_model}",
        "quarter": "{quarter}"
    }},
    "abilities": {{
        "technical_innovation": 评分(1-5),
        "business_impact": 评分(1-5),
        "teamwork": 评分(1-5),
        "project_management": 评分(1-5),
        "cost_awareness": 评分(1-5),
        "strategic_thinking": 评分(1-5)
    }},
    "strengths": [
        "具体优势描述1",
        "具体优势描述2",
        "具体优势描述3",
        "具体优势描述4"
    ],
    "weaknesses": [
        "具体待发展领域1",
        "具体待发展领域2",
        "具体待发展领域3",
        "具体待发展领域4"
    ],
    "growth_suggestions": [
        "具体成长建议1",
        "具体成长建议2",
        "具体成长建议3",
        "具体成长建议4",
        "具体成长建议5",
        "具体成长建议6"
    ],
    "manager_suggestions": [
        "给管理者的建议1",
        "给管理者的建议2",
        "给管理者的建议3",
        "给管理者的建议4",
        "给管理者的建议5",
        "给管理者的建议6"
    ]
}}

## 注意事项
1. 评分要基于述职报告的具体内容，客观公正
2. 优势和待发展领域要具体明确，避免泛泛而谈
3. 成长建议要具有可操作性和针对性
4. 管理者建议要从管理角度提供支持和指导方向
5. 严格按照JSON格式输出，不要添加其他文字说明
"""
    
    try:
        response = dashscope.Generation.call(
            model='qwen-plus',
            prompt=prompt,
            result_format='text',
            api_key=DASHSCOPE_API_KEY
        )
        
        if response.status_code == HTTPStatus.OK:
            response_text = response.output.text.strip()
            # 尝试解析JSON响应
            try:
                # 提取JSON部分（去除可能的markdown格式）
                if '```json' in response_text:
                    json_start = response_text.find('```json') + 7
                    json_end = response_text.find('```', json_start)
                    response_text = response_text[json_start:json_end].strip()
                elif '```' in response_text:
                    json_start = response_text.find('```') + 3
                    json_end = response_text.rfind('```')
                    response_text = response_text[json_start:json_end].strip()
                
                diagnosis_data = json.loads(response_text)
                return diagnosis_data
            except json.JSONDecodeError as e:
                print(f"JSON解析错误: {e}")
                print(f"原始响应: {response_text[:500]}...")
                return None
        else:
            print(f"千问API调用失败: {response.status_code}, {response.message}")
            return None
            
    except Exception as e:
        print(f"千问模型调用异常: {str(e)}")
        return None

def get_fallback_diagnosis(employee_name, ability_model, quarter):
    """
    备用诊断分析结果
    """
    return {
        'employee_info': {
            'name': employee_name,
            'position': ability_model,
            'quarter': quarter
        },
        'abilities': {
            'technical_innovation': 4.0,
            'business_impact': 3.5,
            'teamwork': 4.2,
            'project_management': 3.8,
            'cost_awareness': 3.2,
            'strategic_thinking': 3.6
        },
        'strengths': [
            '工作执行力强，能够按时完成分配的任务',
            '学习能力较好，能够快速掌握新的工作技能',
            '团队合作意识强，与同事关系融洽',
            '工作态度积极，责任心强'
        ],
        'weaknesses': [
            '业务理解深度有待提升，需要更深入了解业务价值',
            '主动性需要加强，可以更积极地提出改进建议',
            '跨部门沟通协调能力需要进一步发展',
            '战略思维和全局观念需要培养'
        ],
        'growth_suggestions': [
            '主动参与业务讨论，深入理解业务需求和价值',
            '定期总结工作经验，形成最佳实践并分享给团队',
            '寻求跨部门合作机会，拓宽工作视野',
            '制定个人发展计划，设定明确的学习和成长目标',
            '参加相关培训课程，提升专业技能和软技能',
            '建立定期反思机制，持续改进工作方法和效率'
        ],
        'manager_suggestions': [
            '为员工提供更多业务培训和学习机会',
            '定期进行一对一沟通，了解员工发展需求',
            '给予员工更多挑战性任务，促进能力提升',
            '建立明确的职业发展路径和晋升标准',
            '提供及时的工作反馈和指导',
            '创造跨部门协作机会，帮助员工拓展视野'
        ]
    }

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)